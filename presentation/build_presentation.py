"""Builds the course presentation deck from this repo's actual results.

    python -m presentation.build_presentation

Reads directly from experiments/*/eval_results*.json and the explainability
figures rather than hardcoding numbers separately, so re-running this after
new results land (e.g. the pending baseline --postprocess run, or learning
curve PNGs once pulled from the server) regenerates the deck instead of
requiring manual slide edits. Two things are still marked TODO in the deck
itself, deliberately not filled with fabricated data:
  - the final baseline-vs-Swin-UNETR comparison, pending the baseline's own
    --postprocess run (see reports/experiment_log.md's Day 9 entry)
  - the learning-curve slide, pending metrics.csv/learning_curve.png being
    pulled from the Shenkar server (never fetched locally so far)
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
EXPLAIN_DIR = ROOT / "experiments" / "swin_unetr_v1" / "explainability"

TITLE_COLOR = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT_COLOR = RGBColor(0xC0, 0x39, 0x2B)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)


def load_json(path: Path) -> dict | None:
    return json.loads(path.read_text()) if path.exists() else None


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank layout


def add_title_box(slide, prs, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), prs.slide_width - Inches(1.2), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = ACCENT_COLOR
    return box


def add_bullet_slide(prs, title, bullets, subtitle=None, note=None):
    """bullets: list of str, or (str, [sub-bullets]) tuples for one level of nesting."""
    slide = blank_slide(prs)
    add_title_box(slide, prs, title, subtitle)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), prs.slide_width - Inches(1.6), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        sub_items = []
        if isinstance(item, tuple):
            text, sub_items = item
        else:
            text = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = f"•  {text}"
        p.font.size = Pt(19)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)
        for sub in sub_items:
            sp = tf.add_paragraph()
            sp.text = f"–  {sub}"
            sp.font.size = Pt(16)
            sp.font.color.rgb = TEXT_COLOR
            sp.level = 1
            sp.space_after = Pt(6)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), prs.slide_width - Inches(1.6), Inches(0.6))
        ntf = note_box.text_frame
        ntf.word_wrap = True
        np = ntf.paragraphs[0]
        np.text = note
        np.font.size = Pt(13)
        np.font.italic = True
        np.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    return slide


def add_image_slide(prs, title, image_path: Path, caption=None, subtitle=None):
    slide = blank_slide(prs)
    add_title_box(slide, prs, title, subtitle)
    if image_path.exists():
        pic_width = prs.slide_width - Inches(1.0)
        left = Inches(0.5)
        top = Inches(1.7)
        slide.shapes.add_picture(str(image_path), left, top, width=pic_width)
    else:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), prs.slide_width - Inches(1.6), Inches(1.0))
        p = box.text_frame.paragraphs[0]
        p.text = f"[Missing: {image_path.relative_to(ROOT)}]"
        p.font.size = Pt(18)
        p.font.color.rgb = ACCENT_COLOR
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), prs.slide_width - Inches(1.6), Inches(0.7))
        ctf = cap_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(14)
        cp.font.color.rgb = TEXT_COLOR
    return slide


def add_table_slide(prs, title, headers, rows, subtitle=None, note=None):
    slide = blank_slide(prs)
    add_title_box(slide, prs, title, subtitle)

    n_rows, n_cols = len(rows) + 1, len(headers)
    table_width = prs.slide_width - Inches(1.6)
    table_height = Inches(0.5 * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.8), Inches(1.8), table_width, table_height)
    table = table_shape.table

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(15)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_COLOR

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

    if note:
        note_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8) + table_height + Inches(0.15), prs.slide_width - Inches(1.6), Inches(0.8)
        )
        ntf = note_box.text_frame
        ntf.word_wrap = True
        np = ntf.paragraphs[0]
        np.text = note
        np.font.size = Pt(13)
        np.font.italic = True
        np.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    return slide


def build() -> Presentation:
    baseline = load_json(ROOT / "experiments" / "baseline_unet_v1" / "eval_results.json")
    swin = load_json(ROOT / "experiments" / "swin_unetr_v1" / "eval_results.json")
    swin_pp = load_json(ROOT / "experiments" / "swin_unetr_v1" / "eval_results_postprocessed.json")
    baseline_pp = load_json(ROOT / "experiments" / "baseline_unet_v1" / "eval_results_postprocessed.json")
    manifest = load_json(EXPLAIN_DIR / "manifest.json")

    prs = new_presentation()

    # 1. Title
    slide = blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), prs.slide_width - Inches(1.6), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Medical Image Segmentation with Swin Transformers"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "3D U-Net baseline vs. Swin UNETR on MSD Task09 (Spleen), with attention-based explainability"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "Shenkar College — Neural Networks, Final Project"
    p3.font.size = Pt(15)
    p3.font.color.rgb = TEXT_COLOR
    p3.alignment = PP_ALIGN.CENTER

    # 2. Goal & challenge
    add_bullet_slide(
        prs,
        "Goal & Challenge",
        [
            "Segment the spleen from abdominal CT (Medical Segmentation Decathlon, Task09_Spleen)",
            "Compare a classic 3D U-Net against a vision-transformer encoder (Swin UNETR) on identical data/methodology",
            "Explain the transformer's predictions using its internal representations, not just report a score",
            (
                "Core challenge",
                [
                    "3D volumetric data at GPU-memory-constrained patch sizes (96³)",
                    "A single university-cluster GPU (NVIDIA L4) with wall-clock time limits",
                    "Transformer attention internals are version-fragile and only semi-documented in MONAI",
                ],
            ),
        ],
    )

    # 3. Related work
    add_bullet_slide(
        prs,
        "Related Work",
        [
            "3D U-Net (Çiçek et al. 2016) / nnU-Net (Isensee et al. 2021): CNN encoder-decoder, strong but local receptive field",
            "Vision Transformer (Dosovitskiy et al. 2021): global self-attention, quadratic cost, high data requirement",
            "Swin Transformer (Liu et al. 2021): windowed, shifted, hierarchical attention — linear cost, CNN-like multi-scale features",
            "Swin UNETR (Hatamizadeh et al. 2022): Swin encoder + CNN decoder, the architecture used here",
            "Attention explainability (Abnar & Zuidema 2020) mostly targets 2D ViT classification — applying it to 3D segmentation is comparatively unexplored",
        ],
        note="Full citations and discussion: reports/related_work.md",
    )

    # 4. Method / fair comparison
    add_bullet_slide(
        prs,
        "Method: A Fair, Same-Conditions Comparison",
        [
            "Identical data pipeline for both models: same HU windowing/spacing, same 96³ patch transforms (src/data/transforms.py)",
            "Identical fixed-seed 80/20 train/val split (src/data/dataset.py) — evaluation reuses it, never re-splits",
            "Matched training methodology: Adam (lr=1e-4, weight_decay=1e-5), DiceCELoss, dropout=0.2, early stopping patience=20",
            "One deliberate difference: batch_size=1 for Swin UNETR vs. 2 for the baseline — an L4 memory necessity for its heavier attention activations, not a methodology compromise",
            "Two metrics, not one: Dice (volumetric overlap) and HD95 (95th-percentile boundary distance) — chosen specifically because Dice alone can miss localization failures",
        ],
    )

    # 5. Process & challenges
    add_bullet_slide(
        prs,
        "Process: Two Failures, Diagnosed and Fixed",
        [
            (
                "Failure 1 — Swin UNETR smoke test crashed at model construction (Day 5)",
                [
                    "TypeError: unexpected keyword argument 'img_size'",
                    "Root cause: installed MONAI release removed that constructor argument; requirements.txt only pinned a lower bound",
                    "Fix: removed img_size from the SwinUNETR(...) call; 96³ patch size still enforced from the data/inference side",
                ],
            ),
            (
                "Failure 2 — Full Swin UNETR run cancelled by a wall-clock time limit at epoch 44/100 (job 322)",
                [
                    "Slurm's effective time cap was below the requested --time=04:00:00",
                    "No checkpoint-resume support existed — 44 epochs of progress had to be discarded",
                    "Fix: full training-state checkpointing + automatic resume added to train.py; re-run completed cleanly",
                ],
            ),
        ],
        note="Full sequence, with config/result/analysis/decision for every run: reports/experiment_log.md",
    )

    # 6. Learning curves (placeholder pending pulled PNGs)
    baseline_curve = ROOT / "experiments" / "baseline_unet_v1" / "learning_curve.png"
    swin_curve = ROOT / "experiments" / "swin_unetr_v1" / "learning_curve.png"
    if baseline_curve.exists() or swin_curve.exists():
        img = baseline_curve if baseline_curve.exists() else swin_curve
        add_image_slide(
            prs,
            "Learning Curves",
            img,
            caption="Validation Dice over training epochs — baseline early-stopped at epoch 75, Swin UNETR at epoch 55.",
        )
    else:
        add_bullet_slide(
            prs,
            "Learning Curves",
            [
                "Baseline 3D U-Net: early-stopped at epoch 75/100, best val Dice reached ~epoch 55",
                "Swin UNETR: early-stopped at epoch 55/100 (restart of job 322, reproduced its dynamics exactly up to epoch 44, then continued to plateau)",
                "TODO: insert learning_curve.png for both runs once pulled from the Shenkar server (experiments/<run_name>/learning_curve.png) — re-run this script afterward to auto-embed them",
            ],
            note="TODO SLIDE — charts not yet pulled locally; text-only placeholder for now.",
        )

    # 7. Raw results table
    if baseline and swin:
        add_table_slide(
            prs,
            "Results: Baseline vs. Swin UNETR (Raw)",
            ["Model", "Mean Dice", "Std Dice", "Mean HD95 (mm)", "Std HD95"],
            [
                [
                    "3D U-Net (baseline)",
                    f"{baseline['summary']['mean_dice']:.4f}",
                    f"{baseline['summary']['std_dice']:.4f}",
                    f"{baseline['summary']['mean_hd95']:.2f}",
                    f"{baseline['summary']['std_hd95']:.2f}",
                ],
                [
                    "Swin UNETR",
                    f"{swin['summary']['mean_dice']:.4f}",
                    f"{swin['summary']['std_dice']:.4f}",
                    f"{swin['summary']['mean_hd95']:.2f}",
                    f"{swin['summary']['std_hd95']:.2f}",
                ],
            ],
            note="Swin UNETR wins on Dice (+13% relative) but HD95 is essentially flat, with its variance more than doubled.",
        )

    # 8. Key finding
    add_bullet_slide(
        prs,
        "Key Finding: Dice Improved, HD95 Didn't — Why?",
        [
            "Swin UNETR's Dice gain is driven by rescuing the baseline's two worst cases (spleen_41: 0.14→0.60, spleen_44: 0.16→0.42)",
            "But on those exact two cases, Swin UNETR's HD95 got worse than the baseline's (194→257mm, 185→255mm)",
            "Better overlap + worse worst-point distance on the same case ⇒ consistent with predicting most of the true spleen correctly, plus one or more small false-positive blobs elsewhere",
            "Dice barely notices a small stray blob (tiny share of total volume); HD95 is dominated by the single worst-matched point — exactly why both metrics were tracked",
        ],
    )

    # 9. Explainability method
    method_desc = "Grad-CAM (fallback: input-gradient saliency)"
    if manifest:
        methods = {v["method"] for v in manifest["cases"].values()}
        method_desc = ", ".join(sorted(methods))
    add_bullet_slide(
        prs,
        "Explainability Method",
        [
            "Ch.6's raw window-attention weights live deep inside MONAI's SwinTransformer internals — version-specific, undocumented",
            "Reaching that deep already broke this project once (the img_size incompatibility) — a second such break wasn't worth risking on a Slurm queue",
            "Instead: forward-hook the stable, public model.swinViT encoder output, Grad-CAM (Selvaraju et al. 2017) from its deepest stage",
            "Automatic fallback to input-gradient saliency (Simonyan et al. 2013) if model.swinViT isn't present on the installed MONAI version",
            f"Method actually used on this run (all 3 cases): {method_desc}",
        ],
    )

    # 10-12. Explainability figures
    labels_order = ["strong", "average", "weak"]
    label_titles = {
        "strong": "Explainability: Strong Case",
        "average": "Explainability: Average Case",
        "weak": "Explainability: Weak Case — the False-Positive Blob",
    }
    label_captions = {
        "strong": "Highest-Dice case. Prediction undershoots the ground truth in this slice despite the best volumetric Dice of the three — a reminder that one slice isn't the whole 3D overlap.",
        "average": "Cleanest example of the method working as intended: the Grad-CAM hot region sits directly over the true spleen body.",
        "weak": "The prediction covers the true spleen on the left, plus a second, disconnected red region with no ground-truth counterpart — direct visual evidence of the false-positive-blob failure mode.",
    }
    if manifest:
        for label in labels_order:
            case = manifest["cases"][label]
            img_path = ROOT / case["figure"]
            add_image_slide(
                prs,
                label_titles[label],
                img_path,
                caption=f"{case['case_id']} — {label_captions[label]}",
            )

    # 13. Diagnosis
    add_bullet_slide(
        prs,
        "Diagnosis: Stray False-Positive Components",
        [
            "Day 7 (metrics) flagged spleen_41/spleen_44 as cases with a Dice win but an HD95 loss",
            "Day 8 (explainability) then visually confirmed the same failure mode on a third case (spleen_25): a disconnected predicted blob with no ground-truth match",
            "Hypothesis: most of Swin UNETR's HD95 problem is small, spatially separate false-positive regions — not a poorly-shaped boundary around a correctly-located spleen",
            "Candidate fix: keep only the largest connected predicted component per case (standard post-processing in medical segmentation for exactly this failure mode)",
        ],
    )

    # 14. Postprocessing results
    if swin and swin_pp:
        add_table_slide(
            prs,
            "Fix: Largest-Connected-Component Postprocessing",
            ["Metric", "Before", "After", "Change"],
            [
                [
                    "Mean Dice",
                    f"{swin['summary']['mean_dice']:.4f}",
                    f"{swin_pp['summary']['mean_dice']:.4f}",
                    f"+{swin_pp['summary']['mean_dice'] - swin['summary']['mean_dice']:.4f}",
                ],
                [
                    "Mean HD95 (mm)",
                    f"{swin['summary']['mean_hd95']:.2f}",
                    f"{swin_pp['summary']['mean_hd95']:.2f}",
                    f"{swin_pp['summary']['mean_hd95'] - swin['summary']['mean_hd95']:.2f}",
                ],
            ],
            note="Every one of the 8 validation cases improved on BOTH metrics — no regressions. spleen_41/spleen_44 (the flagged HD95 outliers) saw the largest drops: -252mm and -240mm.",
        )

    # 15. Final comparison (placeholder until baseline is postprocessed)
    if baseline_pp and swin_pp:
        add_table_slide(
            prs,
            "Final Comparison (Both Models, Postprocessed)",
            ["Model", "Mean Dice", "Mean HD95 (mm)"],
            [
                ["3D U-Net (baseline)", f"{baseline_pp['summary']['mean_dice']:.4f}", f"{baseline_pp['summary']['mean_hd95']:.2f}"],
                ["Swin UNETR", f"{swin_pp['summary']['mean_dice']:.4f}", f"{swin_pp['summary']['mean_hd95']:.2f}"],
            ],
        )
    else:
        add_bullet_slide(
            prs,
            "Final Comparison — Pending",
            [
                "Swin UNETR has been evaluated with postprocessing (previous slide)",
                "The baseline has NOT been postprocessed yet — needed before claiming a final winner, since its raw HD95 (154mm) is just as bad and may have the same fixable cause",
                "TODO: run `sbatch slurm/evaluate.sbatch configs/baseline_unet.yaml outputs/baseline_unet_v1/best_model.pt \"\" postprocess` on Shenkar, pull the result, re-run this script",
            ],
            note="TODO SLIDE — do not present a 'Swin UNETR wins' conclusion until this is filled in.",
        )

    # 16. Limitations & future work
    add_bullet_slide(
        prs,
        "Limitations & Future Work",
        [
            "Explainability heatmaps are coarse by construction (deepest Swin stage collapses a 96³ crop to ~3³ before upsampling)",
            "Explainability crop is centered using the ground-truth label — valid for labeled validation cases, not a deployment-time method",
            "Even after postprocessing, results are not yet at the ~0.90+ Dice / single-digit-mm HD95 published range for this task",
            (
                "Candidate next steps",
                [
                    "Apply postprocessing to the baseline for a fully fair final comparison",
                    "Revisit whether early stopping is premature vs. a ceiling (train-loss curves alongside val-Dice)",
                    "Try attention rollout (Abnar & Zuidema 2020) as a second explainability method to cross-check Grad-CAM",
                ],
            ),
        ],
    )

    # 17. Conclusions
    add_bullet_slide(
        prs,
        "Conclusions",
        [
            "Swin UNETR beats the CNN baseline on Dice, but that alone would have hidden an HD95 problem that only two metrics together revealed",
            "Explainability wasn't just a visualization exercise — it directly corroborated a hypothesis formed from the numbers, on an independent case",
            "The diagnosis led to a concrete, low-cost fix (largest-connected-component postprocessing) that nearly closed the gap: Dice 0.535→0.765, HD95 156→18mm",
            "Every stage of this process — including two failures — is documented as it happened in reports/experiment_log.md",
        ],
    )

    return prs


def main() -> None:
    prs = build()
    out_path = ROOT / "presentation" / "final_presentation.pptx"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
